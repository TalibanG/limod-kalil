# -*- coding: utf-8 -*-
"""LimodKalil — לימוד קליל: עוזר לימודים אישי מבוסס Claude."""
import io
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import uuid

from flask import Flask, Response, jsonify, request, send_from_directory

FROZEN = getattr(sys, "frozen", False)  # True when packaged by PyInstaller

if FROZEN:
    # bundled read-only files (static/) live in the PyInstaller temp dir
    RESOURCE_DIR = sys._MEIPASS
    # writable data lives next to the .exe the user downloaded
    DATA_DIR = os.path.join(os.path.dirname(sys.executable), "LimodKalil-data")
else:
    RESOURCE_DIR = os.path.dirname(os.path.abspath(__file__))
    DATA_DIR = os.path.join(RESOURCE_DIR, "data")

BASE_DIR = RESOURCE_DIR
COURSES_DIR = os.path.join(DATA_DIR, "courses")
CONFIG_PATH = os.path.join(DATA_DIR, "config.json")
os.makedirs(COURSES_DIR, exist_ok=True)

DEFAULT_MODEL = "claude-opus-4-8"
# distributed builds default to API-key auth; the dev copy defaults to the Max subscription
DEFAULT_BACKEND = "api" if FROZEN else "max"
MAX_MATERIAL_CHARS = 600_000  # ~150K tokens, well inside the 1M context window

app = Flask(__name__, static_folder=os.path.join(RESOURCE_DIR, "static"))


# ---------------------------------------------------------------- storage

def load_json(path, default):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return default


def save_json(path, data):
    tmp = path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    os.replace(tmp, path)


def get_config():
    return load_json(CONFIG_PATH, {"api_key": "", "model": DEFAULT_MODEL})


def course_path(cid):
    if not re.fullmatch(r"[0-9a-f]{32}", cid):
        raise ValueError("bad course id")
    return os.path.join(COURSES_DIR, cid + ".json")


def load_course(cid):
    return load_json(course_path(cid), None)


def save_course(course):
    save_json(course_path(course["id"]), course)


def list_courses():
    out = []
    for name in os.listdir(COURSES_DIR):
        if name.endswith(".json"):
            c = load_json(os.path.join(COURSES_DIR, name), None)
            if c:
                out.append({
                    "id": c["id"], "name": c["name"],
                    "materials": len(c.get("materials", [])),
                    "summaries": len(c.get("summaries", [])),
                    "exams": len(c.get("exams", [])),
                    "flashcards": len(c.get("flashcards", [])),
                })
    out.sort(key=lambda c: c["name"])
    return out


# ---------------------------------------------------------------- claude

NO_KEY_MSG = "חסר מפתח API — פתחו את ההגדרות (⚙) והדביקו מפתח מ-platform.claude.com"
NOT_LOGGED_MSG = ("לא מחוברים למנוי Claude — פתחו את ההגדרות (⚙) ולחצו "
                  "\"התחבר למנוי\", ואז אשרו את ההתחברות בדפדפן.")


def find_claude_cli():
    exe = os.path.join(os.path.expanduser("~"), ".local", "bin", "claude.exe")
    if os.path.exists(exe):
        return exe
    return shutil.which("claude")


def cli_model_alias(model):
    if "sonnet" in model:
        return "sonnet"
    if "haiku" in model:
        return "haiku"
    return "opus"


def clean_env():
    """Subprocess env: detach from any parent Claude session, force subscription auth."""
    env = dict(os.environ)
    for k in ("CLAUDECODE", "CLAUDE_CODE_ENTRYPOINT", "CLAUDE_CODE_SESSION_ID",
              "ANTHROPIC_API_KEY", "ANTHROPIC_AUTH_TOKEN"):
        env.pop(k, None)
    return env


def render_prompt(system, messages):
    """Flatten system blocks + conversation into one prompt for CLI mode."""
    parts = [b["text"] for b in system]
    if len(messages) == 1:
        parts.append(messages[-1]["content"])
    else:
        conv = []
        for m in messages[:-1]:
            who = "הסטודנט" if m["role"] == "user" else "אתה (המורה)"
            conv.append(f"{who}: {m['content']}")
        parts.append("היסטוריית השיחה עד כה:\n\n" + "\n\n".join(conv))
        parts.append("הודעת הסטודנט הנוכחית — ענה עליה:\n" + messages[-1]["content"])
    return "\n\n---\n\n".join(parts)


def stream_claude_cli(prompt):
    """Stream text from Claude Code CLI (uses the user's Claude subscription)."""
    exe = find_claude_cli()
    if not exe:
        raise RuntimeError("Claude Code לא מותקן במחשב — התקינו מ-claude.com/code "
                           "או עברו למצב מפתח API בהגדרות.")
    cfg = get_config()
    cmd = [exe, "-p", "--output-format", "stream-json",
           "--include-partial-messages", "--verbose",
           "--model", cli_model_alias(cfg.get("model", DEFAULT_MODEL))]
    flags = 0x08000000 if os.name == "nt" else 0  # CREATE_NO_WINDOW
    with tempfile.TemporaryFile() as errf:
        proc = subprocess.Popen(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE,
                                stderr=errf, env=clean_env(), cwd=DATA_DIR,
                                creationflags=flags)
        try:
            proc.stdin.write(prompt.encode("utf-8"))
            proc.stdin.close()
            got_delta = False
            result_text = None
            for raw in proc.stdout:
                line = raw.decode("utf-8", errors="replace").strip()
                if not line:
                    continue
                try:
                    ev = json.loads(line)
                except json.JSONDecodeError:
                    continue
                t = ev.get("type")
                if t == "stream_event":
                    e = ev.get("event", {})
                    if (e.get("type") == "content_block_delta"
                            and e.get("delta", {}).get("type") == "text_delta"):
                        got_delta = True
                        yield e["delta"]["text"]
                elif t == "result":
                    if ev.get("is_error"):
                        raise RuntimeError(ev.get("result") or "שגיאה בהרצת Claude")
                    result_text = ev.get("result") or ""
            code = proc.wait()
            if code != 0:
                errf.seek(0)
                err = errf.read().decode("utf-8", errors="replace")[-500:]
                if "log" in err.lower() and "/login" in err:
                    raise RuntimeError(NOT_LOGGED_MSG)
                raise RuntimeError("שגיאת Claude Code: " + (err.strip() or f"קוד {code}"))
            if not got_delta and result_text:
                yield result_text
        finally:
            if proc.poll() is None:
                proc.kill()


def parse_json_loose(text):
    text = text.strip()
    m = re.search(r"```(?:json)?\s*(.*?)```", text, re.S)
    if m:
        text = m.group(1).strip()
    start, end = text.find("{"), text.rfind("}")
    if start >= 0 and end > start:
        text = text[start:end + 1]
    return json.loads(text)


def get_client():
    cfg = get_config()
    key = cfg.get("api_key") or os.environ.get("ANTHROPIC_API_KEY", "")
    if not key:
        return None
    import anthropic
    return anthropic.Anthropic(api_key=key)


def materials_block(course, selected_ids=None):
    """Concatenate course materials into one cacheable text block."""
    parts = []
    for m in course.get("materials", []):
        if selected_ids and m["id"] not in selected_ids:
            continue
        parts.append(f"### חומר לימוד: {m['name']}\n\n{m['text']}")
    text = "\n\n---\n\n".join(parts)
    if len(text) > MAX_MATERIAL_CHARS:
        text = text[:MAX_MATERIAL_CHARS] + "\n\n[קוצר — חומר ארוך מדי]"
    return text


def system_blocks(course, role_prompt, selected_ids=None):
    """System prompt: stable role text + course materials with a cache breakpoint."""
    blocks = [{"type": "text", "text": role_prompt}]
    mats = materials_block(course, selected_ids)
    if mats:
        blocks.append({
            "type": "text",
            "text": f"חומרי הקורס \"{course['name']}\":\n\n{mats}",
            "cache_control": {"type": "ephemeral"},
        })
    return blocks


def sse(event_gen):
    """Wrap a text-chunk generator as an SSE response."""
    def stream():
        try:
            for chunk in event_gen:
                yield "data: " + json.dumps({"text": chunk}, ensure_ascii=False) + "\n\n"
            yield "data: " + json.dumps({"done": True}) + "\n\n"
        except Exception as e:  # surface API errors to the UI instead of a dead stream
            yield "data: " + json.dumps({"error": str(e)}, ensure_ascii=False) + "\n\n"
    return Response(stream(), mimetype="text/event-stream",
                    headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


def stream_claude(system, messages, max_tokens=64000):
    cfg = get_config()
    if cfg.get("backend", DEFAULT_BACKEND) == "max":
        yield from stream_claude_cli(render_prompt(system, messages))
    else:
        client = get_client()
        if client is None:
            raise RuntimeError(NO_KEY_MSG)
        with client.messages.stream(
            model=cfg.get("model", DEFAULT_MODEL),
            max_tokens=max_tokens,
            thinking={"type": "adaptive"},
            system=system,
            messages=messages,
        ) as stream:
            yield from stream.text_stream

def stream_and_save(system, messages, on_text):
    """Stream text; ALWAYS persist whatever arrived — even if the stream
    is cut mid-way (window closed, app restarted, network error)."""
    full = []
    completed = False
    try:
        for text in stream_claude(system, messages):
            full.append(text)
            yield text
        completed = True
    finally:
        text = "".join(full).strip()
        if text:
            if not completed:
                text += "\n\n> ⚠️ *התשובה נקטעה באמצע*"
            try:
                on_text(text)
            except Exception:
                pass


def claude_json(system, messages, schema, max_tokens=64000):
    """Non-interactive structured call — returns a validated JSON object."""
    cfg = get_config()
    if cfg.get("backend", DEFAULT_BACKEND) == "max":
        prompt = render_prompt(system, messages) + (
            "\n\n---\n\nהחזר אך ורק JSON תקין (בלי שום טקסט נוסף ובלי code fences) "
            "התואם בדיוק לסכמה הבאה:\n" + json.dumps(schema, ensure_ascii=False)
        )
        text = "".join(stream_claude_cli(prompt))
        return parse_json_loose(text)
    client = get_client()
    if client is None:
        raise RuntimeError(NO_KEY_MSG)
    with client.messages.stream(
        model=cfg.get("model", DEFAULT_MODEL),
        max_tokens=max_tokens,
        thinking={"type": "adaptive"},
        system=system,
        messages=messages,
        output_config={"format": {"type": "json_schema", "schema": schema}},
    ) as stream:
        msg = stream.get_final_message()
    text = next(b.text for b in msg.content if b.type == "text")
    return json.loads(text)


# ---------------------------------------------------------------- file extraction

def extract_text(filename, data):
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                raise RuntimeError("ה-PDF מוגן בסיסמה — הסירו את ההגנה ונסו שוב")
        text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
        if not text.strip():
            raise RuntimeError("לא נמצא טקסט ב-PDF — כנראה קובץ סרוק (תמונות). "
                               "נסו להמיר אותו עם OCR או הדביקו את הטקסט ידנית")
        return text
    if ext == ".docx":
        import docx
        doc = docx.Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    if ext == ".doc":
        raise RuntimeError("קבצי .doc ישנים לא נתמכים — פתחו ב-Word, "
                           "שמרו בתור .docx ונסו שוב")
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- שקופית {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                parts.append("הערות המרצה: " + slide.notes_slide.notes_text_frame.text)
        return "\n".join(parts)
    if ext == ".ppt":
        raise RuntimeError("קבצי .ppt ישנים לא נתמכים — פתחו ב-PowerPoint, "
                           "שמרו בתור .pptx ונסו שוב")
    # txt / md / anything else — try utf-8 then cp1255 (Hebrew Windows)
    for enc in ("utf-8", "cp1255", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


# ---------------------------------------------------------------- routes: app + config

@app.get("/")
def index():
    return send_from_directory(app.static_folder, "index.html")


@app.get("/api/config")
def api_get_config():
    cfg = get_config()
    key = cfg.get("api_key", "")
    return jsonify({
        "backend": cfg.get("backend", DEFAULT_BACKEND),
        "cli_ok": bool(find_claude_cli()),
        "has_key": bool(key or os.environ.get("ANTHROPIC_API_KEY")),
        "key_hint": (key[:11] + "…" + key[-4:]) if len(key) > 18 else "",
        "model": cfg.get("model", DEFAULT_MODEL),
        "frozen": FROZEN,  # distributed build → hide the subscription option
    })


@app.post("/api/config")
def api_set_config():
    body = request.get_json(force=True)
    cfg = get_config()
    if body.get("backend") in ("max", "api"):
        cfg["backend"] = body["backend"]
    if body.get("api_key"):
        cfg["api_key"] = body["api_key"].strip()
    if body.get("model"):
        cfg["model"] = body["model"].strip()
    save_json(CONFIG_PATH, cfg)
    return jsonify({"ok": True})


@app.post("/api/config/test")
def api_test_key():
    try:
        cfg = get_config()
        if cfg.get("backend", DEFAULT_BACKEND) == "max":
            text = "".join(stream_claude_cli("השב במילה אחת בלבד: תקין"))
            if not text.strip():
                raise RuntimeError("לא התקבלה תשובה")
            return jsonify({"ok": True})
        client = get_client()
        if client is None:
            return jsonify({"ok": False, "error": "לא הוגדר מפתח"}), 400
        client.messages.create(
            model=cfg.get("model", DEFAULT_MODEL),
            max_tokens=16,
            messages=[{"role": "user", "content": "שלום"}],
        )
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 400


@app.post("/api/login")
def api_login():
    """Open a console window for the one-time Claude subscription login."""
    exe = find_claude_cli()
    if not exe:
        return jsonify({"error": "Claude Code לא מותקן במחשב"}), 400
    subprocess.Popen(f'start "Claude Login" cmd /k ""{exe}" /login"',
                     shell=True, env=clean_env(), cwd=os.path.expanduser("~"))
    return jsonify({"ok": True})


# ---------------------------------------------------------------- routes: courses & materials

@app.get("/api/courses")
def api_courses():
    return jsonify(list_courses())


@app.post("/api/courses")
def api_create_course():
    name = (request.get_json(force=True).get("name") or "").strip()
    if not name:
        return jsonify({"error": "חסר שם קורס"}), 400
    course = {"id": uuid.uuid4().hex, "name": name, "materials": [],
              "summaries": [], "solutions": [], "exams": [], "flashcards": [], "chat": []}
    save_course(course)
    return jsonify(course)


@app.get("/api/courses/<cid>")
def api_get_course(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    slim = dict(course)
    slim["materials"] = [{"id": m["id"], "name": m["name"], "chars": len(m["text"])}
                         for m in course["materials"]]
    return jsonify(slim)


@app.delete("/api/courses/<cid>")
def api_delete_course(cid):
    try:
        os.remove(course_path(cid))
    except FileNotFoundError:
        pass
    return jsonify({"ok": True})


@app.post("/api/courses/<cid>/materials")
def api_add_material(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    added, failed = [], []
    if request.files:
        for f in request.files.getlist("files"):
            try:
                text = extract_text(f.filename, f.read())
                if not text.strip():
                    raise RuntimeError("הקובץ ריק או שלא נמצא בו טקסט")
            except Exception as e:
                failed.append({"name": f.filename, "error": str(e)})
                continue
            m = {"id": uuid.uuid4().hex, "name": f.filename, "text": text}
            course["materials"].append(m)
            added.append({"id": m["id"], "name": m["name"], "chars": len(text)})
        save_course(course)
        return jsonify({"added": added, "failed": failed})
    else:
        body = request.get_json(force=True)
        m = {"id": uuid.uuid4().hex,
             "name": (body.get("name") or "טקסט מודבק").strip(),
             "text": body.get("text", "")}
        if not m["text"].strip():
            return jsonify({"error": "אין תוכן"}), 400
        course["materials"].append(m)
        added.append({"id": m["id"], "name": m["name"], "chars": len(m["text"])})
    save_course(course)
    return jsonify({"added": added, "failed": failed})


@app.get("/api/courses/<cid>/materials/<mid>")
def api_get_material(cid, mid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    m = next((m for m in course["materials"] if m["id"] == mid), None)
    if not m:
        return jsonify({"error": "חומר לא נמצא"}), 404
    return jsonify(m)


@app.delete("/api/courses/<cid>/materials/<mid>")
def api_del_material(cid, mid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    course["materials"] = [m for m in course["materials"] if m["id"] != mid]
    save_course(course)
    return jsonify({"ok": True})


@app.delete("/api/courses/<cid>/<kind>/<item_id>")
def api_del_item(cid, kind, item_id):
    if kind not in ("summaries", "solutions", "exams", "flashcards"):
        return jsonify({"error": "bad kind"}), 400
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    course[kind] = [x for x in course[kind] if x["id"] != item_id]
    save_course(course)
    return jsonify({"ok": True})


# ---------------------------------------------------------------- routes: AI features

MATH_RULE = ("נוסחאות מתמטיות כתוב אך ורק ב-LaTeX תקני: inline בתוך $...$ "
             "ונוסחאות במרכז בתוך $$...$$. לעולם אל תכתוב LaTeX בלי סימני $.")

SUMMARY_PROMPT = f"""אתה מרצה מצטיין ומומחה ללמידה אפקטיבית, שמכין סיכומים לסטודנטים בישראל.
כתוב בעברית ברורה (מונחים מקצועיים גם באנגלית בסוגריים). {MATH_RULE}

הסיכום בנוי לשתי מטרות: (א) שהסטודנט באמת *יבין ויזכור* את החומר, (ב) שיהיה נוח להעתיק/לכתוב במחברת.
לכן הקפד על העקרונות הבאים, שמבוססים על מדע הלמידה:
- **תמציתי וצפוף** — מילים פשוטות, משפטים קצרים, בלי מליצות. סטודנט צריך יכול להעתיק את זה למחברת בלי לקצר.
- **היררכיה ברורה** — כותרות לפי נושאים, ותחתן נקודות (bullets) קצרות, לא פסקאות ארוכות.
- **הבנה לפני שינון** — לכל מושג מרכזי תן משפט "במילים פשוטות:" שמסביר את האינטואיציה, ואז את ההגדרה הפורמלית.
- **חיבור בין רעיונות** — כשמושג נשען על קודם, ציין זאת במפורש ("נובע מ...", "בניגוד ל...").
- **דוגמה אחת פתורה** לכל רעיון מרכזי — קצרה, עם הנימוק.

מבנה כל פרק:
## [שם הנושא]
**במילים פשוטות:** המהות באינטואיציה.
**הגדרה / נוסחה:** ההגדרה הפורמלית המדויקת.
- נקודות עיקריות קצרות...
**דוגמה:** דוגמה פתורה קצרה.
**🧠 שאלות לבדיקה עצמית:** 2-3 שאלות קצרות שהסטודנט יכול לנסות לענות מהזיכרון (התשובות לא כאן — זה כדי לתרגל שליפה אקטיבית).
**⚠️ טעויות נפוצות:** מה מפספסים בנושא הזה.

בסוף הכול:
## 📋 טבלת סיכום מהיר
טבלה של מושג ↔ מהות בשורה אחת (אידיאלי לחזרה לפני מבחן).

הסיכום צריך להספיק כדי ללמוד למבחן בלי לחזור לחומר המקורי. השתמש ב-Markdown."""

SOLVE_PROMPT = f"""אתה מתרגל מצטיין שפותר מטלות לסטודנטים — המטרה שהסטודנט יבין לעומק ויצליח במבחן, לא רק יעתיק.
פתור בעברית. {MATH_RULE}
אם יש חומרי קורס — פתור בשיטות ובסימונים שלהם בלבד.

הקפד בדיוק על התבנית הבאה לכל שאלה (זה קריטי לקריאוּת):

## 📌 שאלה N — [כותרת קצרה של מה נדרש]

**🎯 מה נבחן כאן:** שורה אחת — המושג או הכלי שהשאלה בודקת.

**💡 רעיון הפתרון:** 2-3 שורות — איך ניגשים ולמה דווקא ככה.

### הפתרון

**צעד 1 — [שם הצעד]:**
מה עושים, ואז שורת ההסבר: *למה?* ...

**צעד 2 — [שם הצעד]:**
... (וכן הלאה — כל צעד קצר, ממוקד, עם "למה" משלו)

> ### ✅ תשובה סופית
> [התשובה, ברורה ושלמה]

**🔎 בדיקת שפיות:** איך מוודאים בשנייה שהתשובה הגיונית (הצבה חזרה, סדר גודל, מקרה קצה).

**⚠️ הטעות הנפוצה:** הטעות שסטודנטים עושים בדיוק בשאלה כזאת.

**📚 למבחן:** איזו וריאציה של השאלה צפויה במבחן ומה משתנה בפתרון.

---

(חזור על התבנית לכל שאלה. בין שאלות — קו מפריד.)"""

CHAT_PROMPT = f"""אתה מורה פרטי סבלני ומעולה לסטודנט ישראלי. ענה בעברית.
הסבר ברמה של הקורס, השתמש בחומרי הקורס כשהם רלוונטיים, תן דוגמאות,
ושאל שאלת בדיקה קצרה בסוף תשובות ארוכות כדי לוודא הבנה. השתמש ב-Markdown. {MATH_RULE}"""

EXAM_PROMPT = f"""אתה כותב מבחנים אוניברסיטאיים. כתוב בעברית.
חבר שאלות ברמת מבחן סוף סמסטר אמיתי על סמך חומרי הקורס: שילוב של הבנה, יישום וחישוב/ניתוח.
לכל שאלה כתוב גם פתרון מלא מוסבר וגם רמז קצר שלא חושף את הפתרון. {MATH_RULE}"""

FOLLOWUP_PROMPT = f"""אתה מתרגל סבלני ומעולה. הסטודנט קיבל פתרון מלא למטלה,
אבל חלק מהצעדים לא ברורים לו והוא שואל עליהם.
ענה בעברית, בגובה העיניים: הסבר את הצעד הלא-ברור לעומק, פרק אותו לתתי-צעדים,
תן דוגמה פשוטה יותר של אותו רעיון אם זה עוזר, וודא בסוף שהקשר לפתרון המקורי ברור.
השתמש ב-Markdown. {MATH_RULE}"""

GRADE_PROMPT = """אתה בודק מבחנים הוגן ומדויק. בדוק את תשובת הסטודנט מול הפתרון הרשמי.
תן ציון 0-100 לכל שאלה, הסבר מה נכון ומה חסר/שגוי, וטיפ קצר לשיפור. כתוב בעברית."""

FLASHCARDS_PROMPT = f"""אתה מכין כרטיסיות זיכרון (flashcards) לסטודנט. כתוב בעברית.
כל כרטיסייה: צד קדמי = שאלה/מונח קצר וממוקד; צד אחורי = תשובה תמציתית ומדויקת.
כסה את כל המושגים, ההגדרות, הנוסחאות והרעיונות המרכזיים בחומר. {MATH_RULE}"""


@app.post("/api/courses/<cid>/summarize")
def api_summarize(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    body = request.get_json(force=True)
    selected = body.get("material_ids") or None
    focus = (body.get("focus") or "").strip()
    if not course["materials"]:
        return jsonify({"error": "אין חומרי לימוד בקורס — העלו קבצים קודם"}), 400

    ask = "כתוב סיכום מלא של חומרי הקורס."
    if focus:
        ask += f"\nהתמקד במיוחד ב: {focus}"
    title = focus or ("סיכום: " + ", ".join(
        m["name"] for m in course["materials"] if not selected or m["id"] in selected)[:80])

    def save_cb(text):
        c = load_course(cid)
        c["summaries"].append({"id": uuid.uuid4().hex, "title": title,
                               "text": text, "ts": time.time()})
        save_course(c)

    return sse(stream_and_save(system_blocks(course, SUMMARY_PROMPT, selected),
                               [{"role": "user", "content": ask}], save_cb))


IMAGE_EXTS = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
              ".webp": "image/webp", ".gif": "image/gif"}
UPLOADS_DIR = os.path.join(DATA_DIR, "uploads")
os.makedirs(UPLOADS_DIR, exist_ok=True)


@app.post("/api/courses/<cid>/solve")
def api_solve(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404

    images = []   # (abs_path, media_type, raw bytes)
    parts = []
    if request.files:
        assignment = (request.form.get("assignment") or "").strip()
        for f in request.files.getlist("files"):
            ext = os.path.splitext(f.filename)[1].lower()
            data = f.read()
            if ext in IMAGE_EXTS:
                path = os.path.join(UPLOADS_DIR, uuid.uuid4().hex + ext)
                with open(path, "wb") as out:
                    out.write(data)
                images.append((path, IMAGE_EXTS[ext], data))
            else:
                try:
                    parts.append(f"### קובץ מטלה: {f.filename}\n\n"
                                 + extract_text(f.filename, data))
                except Exception as e:
                    return jsonify({"error": f"{f.filename}: {e}"}), 400
    else:
        assignment = (request.get_json(force=True).get("assignment") or "").strip()

    full_text = "\n\n".join(([assignment] if assignment else []) + parts).strip()
    if not full_text and not images:
        return jsonify({"error": "הדביקו את המטלה או העלו קובץ"}), 400

    title = (assignment or (parts[0].splitlines()[0].replace("### קובץ מטלה: ", "")
                            if parts else "מטלה מצולמת"))[:80]

    def save_cb(text):
        c = load_course(cid)
        c["solutions"].append({"id": uuid.uuid4().hex, "title": title,
                               "assignment": full_text or "(מטלה מתמונה)",
                               "text": text, "ts": time.time()})
        save_course(c)

    cfg = get_config()
    if cfg.get("backend", DEFAULT_BACKEND) == "max":
        prompt = "המטלה:\n\n" + full_text if full_text else "המטלה מופיעה בתמונות."
        if images:
            paths = "\n".join(p for p, _, _ in images)
            prompt += ("\n\nמצורפות תמונות של המטלה. קרא אותן עם כלי Read "
                       "ופתור את כל השאלות שמופיעות בהן:\n" + paths)
        messages = [{"role": "user", "content": prompt}]
    else:
        import base64
        content = [{"type": "image",
                    "source": {"type": "base64", "media_type": mt,
                               "data": base64.standard_b64encode(raw).decode()}}
                   for _, mt, raw in images]
        content.append({"type": "text",
                        "text": "המטלה:\n\n" + (full_text or "המטלה מופיעה בתמונות המצורפות. פתור את כל השאלות.")})
        messages = [{"role": "user", "content": content}]

    return sse(stream_and_save(system_blocks(course, SOLVE_PROMPT), messages, save_cb))


@app.post("/api/courses/<cid>/solutions/<sid>/ask")
def api_solution_ask(cid, sid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    sol = next((s for s in course["solutions"] if s["id"] == sid), None)
    if not sol:
        return jsonify({"error": "פתרון לא נמצא"}), 404
    question = (request.get_json(force=True).get("question") or "").strip()
    if not question:
        return jsonify({"error": "כתבו שאלה"}), 400

    # persist the question immediately
    sol.setdefault("followups", []).append({"role": "user", "content": question})
    save_course(course)

    intro = (f"המטלה:\n{sol.get('assignment', '')}\n\n"
             f"הפתרון המלא שקיבלתי:\n{sol['text']}\n\n"
             "יש לי שאלות על דרך הפתרון.")
    messages = ([{"role": "user", "content": intro},
                 {"role": "assistant", "content": "בשמחה! מה לא ברור בפתרון?"}]
                + [{"role": m["role"], "content": m["content"]}
                   for m in sol["followups"]][-12:])

    def save_cb(text):
        c = load_course(cid)
        s = next((s for s in c["solutions"] if s["id"] == sid), None)
        if s is not None:
            s.setdefault("followups", []).append(
                {"role": "assistant", "content": text})
            save_course(c)

    return sse(stream_and_save(system_blocks(course, FOLLOWUP_PROMPT),
                               messages, save_cb))


@app.post("/api/courses/<cid>/chat")
def api_chat(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    question = (request.get_json(force=True).get("message") or "").strip()
    if not question:
        return jsonify({"error": "כתבו שאלה"}), 400

    # persist the student's message IMMEDIATELY — before generation starts
    course.setdefault("chat", []).append({"role": "user", "content": question})
    save_course(course)

    history = [{"role": m["role"], "content": m["content"]}
               for m in course["chat"]][-21:]

    def save_cb(text):
        c = load_course(cid)
        c.setdefault("chat", []).append({"role": "assistant", "content": text})
        save_course(c)

    return sse(stream_and_save(system_blocks(course, CHAT_PROMPT), history, save_cb))


@app.post("/api/courses/<cid>/chat/clear")
def api_chat_clear(cid):
    course = load_course(cid)
    if course:
        course["chat"] = []
        save_course(course)
    return jsonify({"ok": True})


EXAM_SCHEMA = {
    "type": "object",
    "properties": {
        "title": {"type": "string"},
        "questions": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "question": {"type": "string"},
                    "hint": {"type": "string"},
                    "solution": {"type": "string"},
                    "points": {"type": "integer"},
                },
                "required": ["question", "hint", "solution", "points"],
                "additionalProperties": False,
            },
        },
    },
    "required": ["title", "questions"],
    "additionalProperties": False,
}


@app.post("/api/courses/<cid>/exam")
def api_exam(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    if not course["materials"]:
        return jsonify({"error": "אין חומרי לימוד בקורס — העלו קבצים קודם"}), 400
    body = request.get_json(force=True)
    n = max(1, min(int(body.get("count", 6)), 15))
    topic = (body.get("topic") or "").strip()
    ask = f"חבר מבחן עם {n} שאלות. סך הנקודות = 100."
    if topic:
        ask += f"\nנושא המבחן: {topic}"
    try:
        exam = claude_json(system_blocks(course, EXAM_PROMPT),
                           [{"role": "user", "content": ask}], EXAM_SCHEMA)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    exam["id"] = uuid.uuid4().hex
    exam["ts"] = time.time()
    course["exams"].append(exam)
    save_course(course)
    return jsonify(exam)


GRADE_SCHEMA = {
    "type": "object",
    "properties": {
        "results": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "score": {"type": "integer"},
                    "feedback": {"type": "string"},
                },
                "required": ["score", "feedback"],
                "additionalProperties": False,
            },
        },
        "total": {"type": "integer"},
        "overall_feedback": {"type": "string"},
    },
    "required": ["results", "total", "overall_feedback"],
    "additionalProperties": False,
}


@app.post("/api/courses/<cid>/exam/grade")
def api_grade(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    body = request.get_json(force=True)
    exam = next((e for e in course["exams"] if e["id"] == body.get("exam_id")), None)
    if not exam:
        return jsonify({"error": "מבחן לא נמצא"}), 404
    answers = body.get("answers", [])
    qa = []
    for i, q in enumerate(exam["questions"]):
        ans = answers[i] if i < len(answers) else ""
        qa.append(f"## שאלה {i+1} ({q['points']} נק')\n{q['question']}\n\n"
                  f"### פתרון רשמי\n{q['solution']}\n\n"
                  f"### תשובת הסטודנט\n{ans or '(לא ענה)'}")
    prompt = ("בדוק את המבחן. לכל שאלה תן score (0-100 יחסית לשאלה) ו-feedback. "
              "total = ציון סופי משוקלל לפי נקודות (0-100).\n\n" + "\n\n---\n\n".join(qa))
    try:
        result = claude_json([{"type": "text", "text": GRADE_PROMPT}],
                             [{"role": "user", "content": prompt}], GRADE_SCHEMA)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    return jsonify(result)


FLASHCARDS_SCHEMA = {
    "type": "object",
    "properties": {
        "cards": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "front": {"type": "string"},
                    "back": {"type": "string"},
                },
                "required": ["front", "back"],
                "additionalProperties": False,
            },
        },
    },
    "required": ["cards"],
    "additionalProperties": False,
}


@app.post("/api/courses/<cid>/flashcards")
def api_flashcards(cid):
    course = load_course(cid)
    if not course:
        return jsonify({"error": "קורס לא נמצא"}), 404
    if not course["materials"]:
        return jsonify({"error": "אין חומרי לימוד בקורס — העלו קבצים קודם"}), 400
    body = request.get_json(force=True)
    n = max(5, min(int(body.get("count", 25)), 60))
    topic = (body.get("topic") or "").strip()
    ask = f"הכן {n} כרטיסיות זיכרון על חומרי הקורס."
    if topic:
        ask += f"\nהתמקד בנושא: {topic}"
    try:
        result = claude_json(system_blocks(course, FLASHCARDS_PROMPT),
                             [{"role": "user", "content": ask}], FLASHCARDS_SCHEMA)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    deck = {"id": uuid.uuid4().hex, "title": topic or "כל החומר",
            "cards": result["cards"], "ts": time.time()}
    course["flashcards"].append(deck)
    save_course(course)
    return jsonify(deck)


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8765))
    print(f"\n  LimodKalil פועל בכתובת:  http://localhost:{port}\n")
    app.run(host="127.0.0.1", port=port, debug=False, threaded=True)
